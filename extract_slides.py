import os
import sys
import re
import cv2
import numpy as np
import yt_dlp
from pptx import Presentation
from pptx.util import Inches
import img2pdf
import requests

def clean_youtube_url(url):
    """URL से फालतू ट्रैकिंग पैरामीटर्स हटाकर शुद्ध URL निकालता है"""
    match = re.search(r'(?:v=|\/)([0-9A-Za-z_-]{11})', url)
    if match:
        video_id = match.group(1)
        return f"https://www.youtube.com/watch?v={video_id}", video_id
    return url, ""

def get_stream_url(youtube_url):
    clean_url, video_id = clean_youtube_url(youtube_url)
    
    # Method 1: yt-dlp with android_vr client (Bypasses Error 152)
    ydl_opts = {
        'format': 'best[height<=720]/best',
        'quiet': False,
        'no_warnings': True,
        'extractor_args': {
            'youtube': {
                'player_client': ['android_vr', 'web_safari', 'android_creator'],
                'player_skip': ['webpage', 'configs']
            }
        }
    }
    
    try:
        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
            info = ydl.extract_info(clean_url, download=False)
            if 'url' in info:
                print("Direct stream URL fetched successfully via android_vr client.")
                return info['url']
    except Exception as e:
        print(f"Direct stream failed: {e}")

    # Method 2: Piped / Invidious API Public Fallbacks
    piped_instances = [
        "https://pipedapi.kavin.rocks",
        "https://api.piped.privacydev.net",
        "https://piped-api.lunar.icu"
    ]
    
    for base in piped_instances:
        try:
            res = requests.get(f"{base}/streams/{video_id}", timeout=10)
            if res.status_code == 200:
                data = res.json()
                streams = data.get('videoStreams', [])
                for s in streams:
                    if s.get('videoOnly') is False or s.get('format') == 'mp4':
                        print(f"Fetched stream from Piped fallback: {base}")
                        return s['url']
        except Exception:
            continue

    raise RuntimeError("Could not retrieve stream URL. All fallback endpoints were throttled.")

def extract_slides(youtube_url, output_folder="slides", sample_interval_sec=5, threshold=0.15):
    os.makedirs(output_folder, exist_ok=True)
    stream_url = get_stream_url(youtube_url)
    
    cap = cv2.VideoCapture(stream_url)
    fps = cap.get(cv2.CAP_PROP_FPS) or 30
    frame_interval = int(fps * sample_interval_sec)
    
    prev_gray = None
    slide_count = 0
    saved_images = []
    frame_idx = 0

    print("Extracting frames and comparing slides...")
    while cap.isOpened():
        cap.set(cv2.CAP_PROP_POS_FRAMES, frame_idx)
        ret, frame = cap.read()
        if not ret:
            break

        gray = cv2.cvtColor(frame, cv2.COLOR_BGR2GRAY)
        gray = cv2.resize(gray, (640, 360))

        if prev_gray is None:
            is_new = True
        else:
            diff = cv2.absdiff(prev_gray, gray)
            non_zero_ratio = np.count_nonzero(diff > 30) / diff.size
            is_new = non_zero_ratio > threshold

        if is_new:
            slide_count += 1
            img_path = f"{output_folder}/slide_{slide_count:04d}.jpg"
            cv2.imwrite(img_path, frame)
            saved_images.append(img_path)
            prev_gray = gray
            print(f"✓ Slide {slide_count} captured at {int(frame_idx/fps)}s")

        frame_idx += frame_interval

    cap.release()
    return saved_images

def save_outputs(images, ppt_path="output.pptx", pdf_path="output.pdf"):
    if not images:
        print("No unique slides found.")
        return

    # PPTX
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for img in images:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(img, 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(ppt_path)

    # PDF
    with open(pdf_path, "wb") as f:
        f.write(img2pdf.convert(images))

    print(f"Generation Complete! Saved {ppt_path} and {pdf_path}")

if __name__ == "__main__":
    url = sys.argv[1] if len(sys.argv) > 1 else "https://www.youtube.com/watch?v=EIhoVK88HOU"
    slides = extract_slides(url)
    save_outputs(slides)
    
